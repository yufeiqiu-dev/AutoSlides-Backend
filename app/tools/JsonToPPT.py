from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
import time
from app.utils.logger import get_logger
logger = get_logger("JSON to PPTX")
def json_to_ppt_bytes(slide_json):
    start = time.time()
    slides = slide_json.get("slides")
    title = slide_json.get("title")
    if not slides or not isinstance(slides, list):
        logger.exception("slide_json must contain a 'slides' list")
        raise ValueError("slide_json must contain a 'slides' list")
    prs = Presentation()
    # Set default slide layout (Title + Content)
    layout = prs.slide_layouts[1]
    if title:
        prs.core_properties.title = title
    for slide_info in slides:
        slide = prs.slides.add_slide(layout)

        # Add title
        title = slide.shapes.title
        title.text = slide_info["header"]

        # Add bullet points
        content = slide.placeholders[1]
        for bullet in slide_info["bullets"]:
            p = content.text_frame.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(18)
    
    # Save to BytesIO instead of a file
    try:
        ppt_bytes = BytesIO()
        prs.save(ppt_bytes)
        ppt_bytes.seek(0)
        elapsed = time.time() - start
        logger.info(f"Successfully parsed json into pptx bytes, time taken: {elapsed:.2f} seconds")
        return ppt_bytes
    except Exception as e:
        raise RuntimeError(f"Failed to generate PowerPoint: {e}")
